import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Custom colors
    BG_COLOR = RGBColor(15, 15, 17)       # Dark grey-black (#0F0F11)
    CARD_BG = RGBColor(26, 26, 30)        # Slightly lighter charcoal (#1A1A1E)
    TEXT_RED = RGBColor(239, 68, 68)      # Neon red (#EF4444)
    TEXT_LIGHT = RGBColor(244, 244, 245)  # Off-white (#F4F4F5)
    TEXT_MUTED = RGBColor(161, 161, 170)  # Gray (#A1A1AA)
    BORDER_COLOR = RGBColor(38, 38, 38)    # Border (#262626)
    
    # Blank layout is usually layout index 6 in default template
    blank_layout = prs.slide_layouts[6]
    
    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header(slide, title_text, category_text="CYBERSATHI HUB PROJECT OVERVIEW"):
        # Header category text
        cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_left = cat_tf.margin_right = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = TEXT_RED
        cat_p.font.name = "Consolas"
        
        # Main header title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.83), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_left = title_tf.margin_right = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_LIGHT
        title_p.font.name = "Arial"

    # ==================== SLIDE 1: COVER ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    
    # Title box (large centered)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "CYBERSATHI HUB"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Arial"
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Multi-Engine Security Telemetry & Orchestration Dashboard"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = "Arial"
    p2.space_before = Pt(15)
    
    # Meta / Info box
    info_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.33), Inches(1.5))
    info_tf = info_box.text_frame
    info_tf.word_wrap = True
    info_tf.margin_left = info_tf.margin_right = info_tf.margin_top = info_tf.margin_bottom = 0
    info_p = info_tf.paragraphs[0]
    info_p.text = "Comprehensive Project Architecture & Component Breakdown"
    info_p.alignment = PP_ALIGN.CENTER
    info_p.font.size = Pt(14)
    info_p.font.bold = True
    info_p.font.color.rgb = TEXT_MUTED
    info_p.font.name = "Consolas"
    
    info_p2 = info_tf.add_paragraph()
    info_p2.text = "Platform: Next.js 16 (React 19) | Supabase SSR | Decoupled Python Heuristics & NLP Engines"
    info_p2.alignment = PP_ALIGN.CENTER
    info_p2.font.size = Pt(11)
    info_p2.font.color.rgb = TEXT_RED
    info_p2.font.name = "Consolas"
    info_p2.space_before = Pt(8)
    
    # ==================== SLIDE 2: ARCHITECTURE ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "System Architecture: Decoupled Microservices")
    
    # Left box: Problem Context
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "THE MONOLITH PROBLEM"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    points = [
      "Dependency Bloat: Running heavy Transformer NLP models along with low-latency TCP socket scanners creates environment collisions (e.g. system packages vs. python version bindings).",
      "Resource Contention: CPU/RAM-heavy analysis (like deep email classification) halts the main UI loop in monolithic frameworks.",
      "Scalability Limits: Decoupling allows isolated engines to scale horizontally, running on separate machines or ports without rebuilding the entire UI."
    ]
    for pt in points:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(12)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: Solution Architecture
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "CYBERSATHI HUB SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    solutions = [
      "Next.js 16 App Router: Serves as the unified command center UI. Compiles optimized static pages for low-latency view rendering.",
      "Secure Server Actions: Next.js acts as a proxy, executing fetch operations to local loopback ports (Port 5000-5008). Avoids CORS blocks and secures backend APIs.",
      "Supabase Cloud Auth: Handles secure operator logins and PostgreSQL telemetry data synchronization with built-in Row-Level Security (RLS)."
    ]
    for sol in solutions:
        ap = right_tf.add_paragraph()
        ap.text = "• " + sol
        ap.font.size = Pt(12)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 3: PORT MAPPING ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Microservices Port Registry & Core Tech Stack")
    
    # Table layout
    rows = 9
    cols = 4
    left = Inches(0.75)
    top = Inches(1.8)
    width = Inches(11.83)
    height = Inches(4.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(1.5) # Port
    table.columns[1].width = Inches(2.8) # Engine Name
    table.columns[2].width = Inches(2.2) # Core Tech Stack
    table.columns[3].width = Inches(5.33) # Functional Summary
    
    headers = ["PORT", "ENGINE MODULE", "TECH STACK", "FUNCTIONAL DESCRIPTION"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEXT_RED
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "Consolas"
        
    data = [
      ("5000", "Password Analyzer", "Python (Flask) + zxcvbn", "Shannon entropy core and breach dictionary lookups."),
      ("5001", "URL Threat Detector", "Python (Flask) + Heuristics", "Heuristic scans & domain redirects verification."),
      ("5002", "Hash Identifier", "Python (Flask) + HashesDB", "MD5, SHA-256 algorithm matching & signatures."),
      ("5003", "Port Scanner", "Python (Flask) + SocketAsync", "Asynchronous TCP SYN scanning for open ports."),
      ("5004", "Malware Sandbox", "Python (Flask) + YARA", "Static signature scanning for trojans & malware."),
      ("5005", "IP Tracker", "Python (Flask) + GeoIP2", "IP geolocating & threat intelligence blacklisting."),
      ("5006", "SSL Inspector", "Python (Flask) + PyOpenSSL", "X.509 cert chains, issuer trust, and expiry auditor."),
      ("5008", "Email Phishing Classifier", "Python (Flask) + DistilBERT", "NLP classification for phishing & social engineering.")
    ]
    
    for r, (port, engine, tech, desc) in enumerate(data):
        row_cells = [port, engine, tech, desc]
        for c, text in enumerate(row_cells):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            if c < 3:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Consolas"
            else:
                p.alignment = PP_ALIGN.LEFT
                p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT

    # ==================== SLIDE 4: PASSWORD & URL ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Engine Detail: Password & URL Evaluators")
    
    # Left box: Password Analyzer
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "PASSWORD STRENGTH EVALUATOR (PORT 5000)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    pass_pts = [
      "Shannon Entropy Core: Computes mathematical randomness on the character distribution of inputs, quantifying the information content.",
      "Zxcvbn Modeling: Mimics human cracking patterns by scanning for keyboard spatial patterns, repeating sequences, and common dates.",
      "Breach Database Checks: Interrogates local indices of cracked password dictionaries to flag historically leaked user keys.",
      "Estimated Crack Time: Provides precise calculation models mapping brute-force times at different scanning speeds."
    ]
    for pt in pass_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: URL Threat Detector
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "URL HEURISTIC THREAT SCANNER (PORT 5001)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    url_pts = [
      "Lexical Analysis: Scans URL string structure for typical indicator traits of fraud (excessive dots, fake subdomains, credential targets).",
      "Entropy Calculations: Computes character entropy on target domain names to detect random-looking domain names typical in botnets.",
      "Active Redirection Check: Uses HTTP client tunnels to follow redirect hops, ensuring safety of final destination sites.",
      "Phishing Blacklists: Cross-references domain names against global repositories mapping malicious sites."
    ]
    for pt in url_pts:
        ap = right_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 5: HASH & PORT SCANNER ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Engine Detail: Cryptographic Identifiers & Port Scanners")
    
    # Left box: Hash Identifier
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "HASH ALGORITHM IDENTIFIER (PORT 5002)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    hash_pts = [
      "Signature Length Verification: Leverages hash byte length and patterns to filter algorithm formats (MD5-32 char, SHA-256-64 char).",
      "Regular Expression Matchers: Scans character ranges to identify special systems (bcrypt, argon2, phpass, django hashes).",
      "Cryptographic Profiling: Helps cyber analysts determine hash types before feeding target digests into cracking grids.",
      "Database Matching: Maps inputs to known hash format patterns for 200+ distinct hashing functions."
    ]
    for pt in hash_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: Port Scanner
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "ASYNCHRONOUS PORT SCANNER (PORT 5003)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    port_pts = [
      "SYN Scanning Emulation: Initiates socket connection sequences to probe for active system communication ports.",
      "Multi-Threaded Probes: Scans target hosts concurrently, drastically shortening analysis windows compared to single-threaded routines.",
      "Common Ports Catalog: Probes target ports (e.g. 21/FTP, 22/SSH, 80/HTTP, 443/HTTPS, 3306/MySQL) to identify vulnerable service entrypoints.",
      "Service Resolution: Interrogates banners returned by target services to trace platform versions."
    ]
    for pt in port_pts:
        ap = right_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 6: MALWARE SIGNATURES ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Engine Detail: Static Malware Sandboxing")
    
    # Left box: Malware Engine
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "STATIC FILE ANALYZER (PORT 5004)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    mal_pts = [
      "YARA Rule Processing: Integrates customizable YARA rule files directly into the scanning sequence.",
      "Signature File Scans: Scans files without local code execution, securing the hosting server from compromise.",
      "Metadata Extractor: Extracts file headers, binary formats (PE, ELF, Mach-O), size metrics, and hashes.",
      "Heuristic Signature Engine: Matches strings and raw hex sequences to identify hidden malware payload footprints."
    ]
    for pt in mal_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: Workflow Diagram
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "MALWARE SCAN PIPELINE FLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    pipeline_pts = [
      "1. Dropzone Upload: Client frontend (React) pushes selected file stream via multipart/form-data payload.",
      "2. Server Action Tunnel: Next.js API intercepts stream, validates headers, and forwards buffer to loopback Port 5004.",
      "3. Engine Evaluation: Flask YARA daemon compiles custom signatures, scans byte-stream, and checks MD5/SHA-256.",
      "4. JSON Threat Profile: Returns telemetry matching triggers, severity rankings, and quarantine recommendations.",
      "5. Dashboard Display: Renders matrix-themed red/green threat results to the operator console."
    ]
    for pt in pipeline_pts:
        ap = right_tf.add_paragraph()
        ap.text = pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 7: IP & SSL ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Engine Detail: IP Threat Trackers & SSL Cert Inspectors")
    
    # Left box: IP Tracker
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "IP GEOLOCATOR & REPUTATION ENGINE (PORT 5005)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    ip_pts = [
      "GeoIP Database Querying: Matches target IP addresses against localized coordinates (latitude, longitude, country).",
      "Network ASN Resolution: Fetches hosting provider and Autonomous System Number (ASN) tracking.",
      "Reputation Checkers: Queries threat databases to check if IPs are linked to botnet activity or proxy networks.",
      "DDoS Origin Mapping: Aids in packet traces, helping operators localize network attack sources."
    ]
    for pt in ip_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: SSL Inspector
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "X.509 SSL/TLS CERTIFICATE INSPECTOR (PORT 5006)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    ssl_pts = [
      "Handshake Parsing: Initiates TLS handshake sequence to extract target SSL X.509 certificate streams.",
      "Cryptographic Audit: Details public key algorithms, key bit strength, and signature algorithm schemas.",
      "Trust Chain Validation: Scans certificate issuing authority (CA) fields to confirm root trust validation.",
      "Expiration Watchdog: Identifies exact valid dates to flag expiring or invalid domain credentials."
    ]
    for pt in ssl_pts:
        ap = right_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 8: AI EMAIL PHISHING ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "AI NLP Modeling: Phishing Mail Classifier (Port 5008)", "AI CORE MODELING DEPTH")
    
    # Left box: NLP Model Architecture
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "DISTILBERT NLP THREAT MODEL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    ai_pts = [
      "Architecture Choice: Uses DistilBERT (a distilled, lightweight version of BERT) maintaining 97% of BERT's performance with a 40% reduction in size.",
      "Attention Mechanism: Utilizes transformer multi-head self-attention layer weights to trace syntactic relationships across input text blocks.",
      "Fine-Tuning: Trained on extensive corpora of scam, spam, phishing, and safe corporate emails.",
      "Why Port Isolation? Transformer runtime loading requires torch and huggingface dependencies. Separating this engine avoids blocking low-latency socket routines."
    ]
    for pt in ai_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: Classification Features
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "CLASSIFICATION PIPELINE & METRICS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    pipeline_pts = [
      "Tokenization Process: Raw email string text is split using WordPiece tokenizers, mapping terms to vocabulary IDs.",
      "Heuristic Context Parsing: Analyzes subject lines and email body headers for sender impersonation and urgent CTAs.",
      "Decision Outputs: Calculates classification probabilities, flagging emails as: SAFE, SUSPICIOUS, or PHISHING.",
      "Token Importance Mapping: Extracts key alert tokens (e.g. 'immediately', 'verify bank', 'login alert') contributing to high threat scores."
    ]
    for pt in pipeline_pts:
        ap = right_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 9: THREAT INTEL ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Global Threat Feed Aggregation Flow")
    
    # Left box: Threat Feed Architecture
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = CARD_BG
    left_shape.line.color.rgb = BORDER_COLOR
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_right = left_tf.margin_top = left_tf.margin_bottom = Inches(0.4)
    
    p = left_tf.paragraphs[0]
    p.text = "XML RSS THREAT INTEGRATOR"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    rss_pts = [
      "CISA Advisories parsing: Connects directly to cisa.gov cybersecurity advisory bulletins.",
      "BleepingComputer RSS parsing: Pulls commercial industry news from BleepingComputer feeds.",
      "Chronological Sorting: Parses XML date metadata, normalizing varying date formats, and sorting feeds descending.",
      "Relative Time Engine: Calculates date diffs, displaying durations (e.g. '4m ago', '3h ago') dynamically on UI."
    ]
    for pt in rss_pts:
        ap = left_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # Right box: Code Flow Logic
    right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = CARD_BG
    right_shape.line.color.rgb = BORDER_COLOR
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_right = right_tf.margin_top = right_tf.margin_bottom = Inches(0.4)
    
    p = right_tf.paragraphs[0]
    p.text = "ROBUST OFFLINE FAILSAFE SYSTEM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_RED
    p.font.name = "Consolas"
    p.space_after = Pt(12)
    
    fail_pts = [
      "Try-Catch Isolation: Wrapped in strict error wrappers. Prevents external connection timeout from throwing React build/runtime errors.",
      "Client Caching Fallback: If external feeds are unreachable (network downtime, CORS, offline environment), the module falls back to local threat intelligence caches.",
      "No Blocking: Next.js server action executes asynchronously, ensuring home dashboard render finishes instantly without delay."
    ]
    for pt in fail_pts:
        ap = right_tf.add_paragraph()
        ap.text = "• " + pt
        ap.font.size = Pt(11)
        ap.font.color.rgb = TEXT_LIGHT
        ap.font.name = "Arial"
        ap.space_after = Pt(8)

    # ==================== SLIDE 10: CORE ENHANCEMENTS ====================
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_header(slide, "Recent Enhancements: Branding, Persistence & UX")
    
    # Grid of three box shapes for 3 key updates
    box_width = Inches(3.68)
    box_height = Inches(4.8)
    spacing = Inches(0.4)
    start_left = Inches(0.75)
    top_pos = Inches(1.8)
    
    updates = [
      ("🎨 BRANDING REFACTOR", 
       "Cybersathi Hub Branding", 
       ["Next.js metadata update: Set document title to 'Cybersathi Hub' for cleaner operator experience.", 
        "Startup sequence: Redesigned the terminal boot sequence logs to output 'INIT: Cybersathi Hub Core'.", 
        "Header UI: Set branding header to 'CYBERSATHI HUB CORE'."]),
      ("🔒 BACK-NAV PERSISTENCE", 
       "Tab-Level Session Storage", 
       ["State Isolation fix: Returning to Command Center from tools was resetting home state to 'login'.", 
        "sessionStorage Hook: Integrated session authentication storage flag on successful logins.", 
        "Autopass filter: Skips login sequence on return mounts within the same tab."]),
      ("⚡ MODERN BOX-CARD UX", 
       "Dynamic UI Component Overhaul", 
       ["Grid Restructuring: Arranged vertical box toolcards in a 3-column responsive layout.", 
        "Neon Hover Animation: Hovering over card triggers red-team borders and bottom radial glow.", 
        "Accent & Tags: Slide-in top header red accent line and absolute-positioned 'READY' tags."])
    ]
    
    for i, (title, subtitle, pts) in enumerate(updates):
        left_pos = start_left + i * (box_width + spacing)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = BORDER_COLOR
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_RED
        p.font.name = "Consolas"
        p.space_after = Pt(4)
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.font.name = "Arial"
        p_sub.space_after = Pt(16)
        
        for pt in pts:
            ap = tf.add_paragraph()
            ap.text = "• " + pt
            ap.font.size = Pt(10.5)
            ap.font.color.rgb = TEXT_LIGHT
            ap.font.name = "Arial"
            ap.space_after = Pt(8)

    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    out = "C:\\Users\\adity\\Downloads\\Sentinel-Dashboard-main\\Cybersathi_Hub_Project_Presentation.pptx"
    if len(sys.argv) > 1:
        out = sys.argv[1]
    create_presentation(out)
